"""
metrics/general_stats_eval.py — General presentation statistics.

Computes structural statistics for each presentation:
  - pages:      Number of slides (from PPTX/PDF, fallback to text headers)
  - characters: Total character count from extracted markdown text
  - words:      Total word count from extracted markdown text
  - figures:    Number of images/figures (from PPTX/PDF, 0 if unavailable)

Character and word counts are computed from the pipeline's extracted markdown
text (processed_data/) rather than raw PPTX/PDF parsing, so they reflect the
same text that all other metrics evaluate.

These are absolute per-method metrics (not pairwise).
No LLM required.

Dependencies: pip install python-pptx pymupdf
Output: results/general_stats_eval.json
"""

import re
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))
import constants as C
from utils.result_utils import (
    load_existing,
    make_metadata,
    read_processed_text,
    result_path,
    save_incremental,
)

METRIC_NAME = "general_stats_eval"


def _count_pptx_structure(pptx_path: Path) -> dict | None:
    """
    Extract page count and figure count from a PPTX file.

    Returns dict with keys: pages, figures
    """
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(str(pptx_path))
        pages = len(prs.slides)
        figures = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    figures += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    if hasattr(shape, "image"):
                        figures += 1

        return {"pages": pages, "figures": figures}
    except ImportError:
        print(f"  [stats] python-pptx not installed. Run: pip install python-pptx")
        return None
    except Exception as e:
        print(f"  [stats] Error reading {pptx_path}: {e}")
        return None


def _count_pdf_structure(pdf_path: Path) -> dict | None:
    """
    Extract page count and figure count from a PDF file.

    Returns dict with keys: pages, figures
    """
    try:
        import pymupdf

        doc = pymupdf.open(str(pdf_path))
        pages = len(doc)
        figures = 0

        for page in doc:
            figures += len(page.get_images(full=True))

        doc.close()
        return {"pages": pages, "figures": figures}
    except ImportError:
        print(f"  [stats] pymupdf not installed. Run: pip install pymupdf")
        return None
    except Exception as e:
        print(f"  [stats] Error reading {pdf_path}: {e}")
        return None


def _count_text_stats(paper_name: str, method: str) -> dict | None:
    """
    Compute character and word counts from the extracted markdown text.
    Optionally compute page count from '## Slide N' headers if no PPTX/PDF
    structure is available.

    Returns dict with keys: characters, words, pages (from text headers)
    """
    text = read_processed_text(paper_name, method)
    if text is None:
        return None

    # Count slides by "## Slide N" headers
    slides = re.findall(r"^## Slide \d+", text, re.MULTILINE)
    pages = len(slides) if slides else 1

    characters = len(text)
    words = len(text.split())

    return {
        "pages": pages,
        "characters": characters,
        "words": words,
    }


def _find_presentation(method: str, paper_name: str) -> Path | None:
    """Find a PPTX or PDF file for a method/paper combination."""
    method_dir = Path(C.GENERATED_SAMPLES_DIR) / method / paper_name
    if not method_dir.exists():
        return None
    # Try PPTX first
    pptx_files = list(method_dir.glob("*.pptx"))
    if pptx_files:
        return pptx_files[0]
    # Try PDF
    pdf_files = list(method_dir.glob("*.pdf"))
    if pdf_files:
        return pdf_files[0]
    return None


def run(papers: list[str], baseline_methods: list[str]) -> dict:
    """
    Compute general statistics for all methods (ours + baselines) on all papers.

    Page count and figure count come from the raw PPTX/PDF files.
    Character count and word count come from the extracted markdown text,
    so they reflect the same text that all other metrics evaluate.
    """
    out_path = result_path(METRIC_NAME)
    existing = load_existing(out_path)
    per_paper: dict = existing.get("per_paper", {})
    metadata = make_metadata(METRIC_NAME, "none (no LLM)")

    all_methods = [C.OURS_METHOD] + baseline_methods

    for i, paper in enumerate(papers, 1):
        print(f"\n[{METRIC_NAME}] [{i}/{len(papers)}] {paper}")

        if paper not in per_paper:
            per_paper[paper] = {}

        for method in all_methods:
            if method in per_paper.get(paper, {}):
                print(f"  Skipping {method} (already done)")
                continue

            # --- Text stats (characters, words) from extracted markdown ---
            text_stats = _count_text_stats(paper, method)
            if text_stats is None:
                print(f"  Skipping {method}: extracted text not found")
                continue

            characters = text_stats["characters"]
            words = text_stats["words"]

            # --- Structural stats (pages, figures) from PPTX/PDF ---
            pres_path = _find_presentation(method, paper)
            structure = None
            if pres_path:
                if pres_path.suffix.lower() == ".pptx":
                    structure = _count_pptx_structure(pres_path)
                else:  # .pdf
                    structure = _count_pdf_structure(pres_path)

            if structure:
                pages = structure["pages"]
                figures = structure["figures"]
            else:
                # Fallback: page count from text headers, no figure count
                pages = text_stats["pages"]
                figures = 0

            stats = {
                "pages": pages,
                "characters": characters,
                "words": words,
                "figures": figures,
            }
            per_paper[paper][method] = stats
            print(
                f"  {method}: pages={stats['pages']}, "
                f"chars={stats['characters']}, words={stats['words']}, "
                f"figures={stats['figures']}"
            )

        save_incremental(out_path, {"metadata": metadata, "per_paper": per_paper})

    # Per-method summary
    per_method: dict[str, dict] = {}
    for method in all_methods:
        pages_list = []
        chars_list = []
        words_list = []
        figs_list = []
        for p in per_paper:
            if method in per_paper[p]:
                s = per_paper[p][method]
                pages_list.append(s["pages"])
                chars_list.append(s["characters"])
                words_list.append(s.get("words", 0))
                figs_list.append(s["figures"])

        n = len(pages_list)
        per_method[method] = {
            "mean_pages": sum(pages_list) / n if n else 0.0,
            "mean_characters": sum(chars_list) / n if n else 0.0,
            "mean_words": sum(words_list) / n if n else 0.0,
            "mean_figures": sum(figs_list) / n if n else 0.0,
            "papers_evaluated": n,
        }

    metadata["total_papers"] = len(per_paper)
    final = {
        "metadata": metadata,
        "per_method_summary": per_method,
        "per_paper": per_paper,
    }
    save_incremental(out_path, final)
    print(f"\n[{METRIC_NAME}] Done. Results -> {out_path}")
    return final
