"""
data_prep/convert_to_text.py — Data preparation step.

Converts:
  • PPTX presentations → markdown text (extracted directly via python-pptx)
  • PDF presentations  → markdown text (extracted via pymupdf4llm)
  • Image-only slides  → markdown text (extracted via GPT-4o vision)
  • Original paper PDFs → markdown text (with supplementary material stripped)

When python-pptx or pymupdf yields empty slides (text is embedded in images),
the pipeline converts slides to PNG images and sends each to GPT-4o vision
for OCR-quality text extraction. This requires OPENAI_API_KEY to be set.

Output folder structure
-----------------------
  PROCESSED_DATA_DIR/
    {paper}/
      orig.md           ← source paper (PDF → markdown, supp material removed)
      {method}.md        ← slide text (PPTX/PDF/GPT-4o → markdown)
      ...

Files are skipped automatically if they already exist (pass force=True to re-convert).
"""

import re
import sys
import time
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))
import constants as C

try:
    import pymupdf4llm
except ImportError:
    pymupdf4llm = None


# ── Supplementary material stripping ─────────────────────────────────────────

# Regex patterns that match reference section headings in markdown
_REFERENCES_HEADING_RE = re.compile(
    r"^(#{1,4})\s+"                                     # 1-4 '#' characters + space
    r"(References|Bibliography|Works\s+Cited|Literature\s+Cited)"
    r"\s*$",
    re.IGNORECASE | re.MULTILINE,
)

# Regex patterns that match supplementary / appendix headings in markdown.
# These are the sections we want to REMOVE when they appear AFTER references.
_SUPPLEMENTARY_HEADING_RE = re.compile(
    r"^(#{1,4})\s+"                                     # 1-4 '#' characters + space
    r"("
    r"Appendix(?:\s+[A-Z0-9])?"                         # Appendix, Appendix A, Appendix 1
    r"|Appendices"
    r"|Supplementary\s+Materials?"
    r"|Supplementary\s+Information"
    r"|Supplemental\s+Materials?"
    r"|Supporting\s+Information"
    r"|Online\s+Appendix"
    r"|Additional\s+Materials?"
    r"|Supplementary"                                    # standalone
    r")"
    r"(\s.*)?$",                                         # optional trailing text on same line
    re.IGNORECASE | re.MULTILINE,
)


def strip_supplementary_material(text: str) -> tuple[str, bool]:
    """
    Remove supplementary material / appendices that appear AFTER the references
    section in an academic paper.

    Logic:
      1. Find the LAST references heading in the document.
      2. Starting from the end of that references heading, scan forward for
         section headings that indicate supplementary / appendix content.
      3. If found, cut everything from that supplementary heading onward.
      4. References themselves are ALWAYS kept intact.

    Returns:
        (processed_text, was_stripped) — the cleaned text and whether anything was cut.

    Safety guarantees:
      - If no references section is found → text returned unchanged.
      - If no supplementary heading is found after references → text returned unchanged.
      - References section content is NEVER removed.
      - Content BEFORE references is NEVER removed.
    """
    ref_matches = list(_REFERENCES_HEADING_RE.finditer(text))
    if not ref_matches:
        return text, False

    ref_match = ref_matches[-1]
    ref_heading_level = len(ref_match.group(1))
    after_ref_start = ref_match.end()

    supp_matches = list(_SUPPLEMENTARY_HEADING_RE.finditer(text, after_ref_start))
    if not supp_matches:
        return text, False

    cut_position = None
    for supp_match in supp_matches:
        supp_heading_level = len(supp_match.group(1))
        if supp_heading_level <= ref_heading_level + 1:
            cut_position = supp_match.start()
            break

    if cut_position is None:
        return text, False

    stripped_text = text[:cut_position].rstrip() + "\n"
    return stripped_text, True


# ── Conversion helpers ────────────────────────────────────────────────────────

def convert_pptx_to_text(pptx_path: Path) -> str | None:
    """
    Extract text directly from a PPTX file using python-pptx.

    Returns a markdown-formatted string with slide headings and content.
    """
    try:
        from pptx import Presentation
    except ImportError:
        print("  [convert] python-pptx is not installed. Run: pip install python-pptx")
        return None

    try:
        prs = Presentation(str(pptx_path))
        parts = []

        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)

            slide_content = "\n".join(slide_texts) if slide_texts else "(empty slide)"
            parts.append(f"## Slide {i}\n{slide_content}")

        return "\n\n".join(parts)
    except Exception as e:
        print(f"  [convert] Error reading PPTX {pptx_path}: {e}")
        return None


def convert_pdf_to_markdown(pdf_path: Path) -> str | None:
    """Convert a PDF to markdown text using pymupdf4llm."""
    if pymupdf4llm is None:
        print("  [convert] pymupdf4llm is not installed. Run: pip install pymupdf4llm")
        return None
    try:
        return pymupdf4llm.to_markdown(str(pdf_path))
    except Exception as e:
        print(f"  [convert] Error converting PDF {pdf_path}: {e}")
        return None


def convert_pdf_slides_to_text(pdf_path: Path) -> str | None:
    """
    Convert a PDF presentation to markdown text using pymupdf4llm.

    Each page is treated as a slide. Returns markdown with slide headings.
    """
    if pymupdf4llm is None:
        print("  [convert] pymupdf4llm is not installed. Run: pip install pymupdf4llm")
        return None
    try:
        import pymupdf
        doc = pymupdf.open(str(pdf_path))
        parts = []
        for i, page in enumerate(doc, 1):
            text = page.get_text().strip()
            if not text:
                text = "(empty slide)"
            parts.append(f"## Slide {i}\n{text}")
        doc.close()
        return "\n\n".join(parts)
    except Exception as e:
        print(f"  [convert] Error converting PDF slides {pdf_path}: {e}")
        return None


# ── GPT-4o vision text extraction ────────────────────────────────────────────

def _load_extraction_prompt() -> str:
    """Load the vision text extraction prompt template."""
    path = Path(C.PROMPTS_DIR) / "vision_extract_text.txt"
    if not path.exists():
        raise FileNotFoundError(f"Extraction prompt not found: {path}")
    return path.read_text(encoding="utf-8")


def extract_text_via_gpt4o(pres_file: Path, method: str, paper_name: str) -> str | None:
    """
    Extract text from slide images using GPT-4o vision.

    Used as a fallback when python-pptx / pymupdf yields empty slides
    (i.e. the text is embedded in images rather than in text frames).

    Steps:
      1. Convert the presentation to PNG images (cached in IMAGES_CACHE_DIR).
      2. Send each slide image to GPT-4o with the extraction prompt.
      3. Combine responses into the standard "## Slide N" markdown format.

    Returns markdown text or None on failure.
    """
    if not C.VISION_EXTRACTION_ENABLED:
        print("    [gpt4o] Vision extraction is disabled (VISION_EXTRACTION_ENABLED=False)")
        return None

    if not C.OPENAI_API_KEY:
        print("    [gpt4o] Skipping: OPENAI_API_KEY is not set")
        return None

    from llm.client import call_vision
    from utils.image_utils import slides_to_images

    # Convert slides to images
    out_dir = Path(C.IMAGES_CACHE_DIR) / method / paper_name
    images = slides_to_images(pres_file, out_dir)
    if not images:
        print("    [gpt4o] No images produced from presentation")
        return None

    prompt = _load_extraction_prompt()
    parts = []
    success_count = 0

    for i, img_path in enumerate(images, 1):
        print(f"    [gpt4o] Extracting text from slide {i}/{len(images)} ...", end=" ", flush=True)
        response = call_vision(
            prompt,
            [img_path],
            model=C.VISION_EXTRACTION_MODEL,
            backend="openai",
            max_tokens=C.MAX_TOKENS,
            temperature=0.0,
        )
        if response and response.strip():
            text = response.strip()
            success_count += 1
        else:
            text = "(empty slide)"
        parts.append(f"## Slide {i}\n{text}")
        print("done")

        if C.SLEEP_BETWEEN_CALLS > 0 and i < len(images):
            time.sleep(C.SLEEP_BETWEEN_CALLS)

    # If no slide was successfully extracted, return None so fallbacks can kick in
    if success_count == 0:
        print("    [gpt4o] All slides failed extraction — returning None")
        return None

    return "\n\n".join(parts)


def _extract_from_existing_images(images: list[str]) -> str | None:
    """Extract text from a list of already-existing slide images via GPT-4o."""
    from llm.client import call_vision

    prompt = _load_extraction_prompt()
    parts = []

    for i, img_path in enumerate(images, 1):
        print(f"    [gpt4o] Extracting text from slide {i}/{len(images)} ...", end=" ", flush=True)
        response = call_vision(
            prompt,
            [img_path],
            model=C.VISION_EXTRACTION_MODEL,
            backend="openai",
            max_tokens=C.MAX_TOKENS,
            temperature=0.0,
        )
        text = response.strip() if response else "(empty slide)"
        parts.append(f"## Slide {i}\n{text}")
        print("done")

        if C.SLEEP_BETWEEN_CALLS > 0 and i < len(images):
            time.sleep(C.SLEEP_BETWEEN_CALLS)

    result = "\n\n".join(parts)
    # Check if all slides are empty
    non_empty = [
        line for line in result.split("\n")
        if line.strip()
        and not line.strip().startswith("## Slide")
        and line.strip() != "(empty slide)"
    ]
    return result if non_empty else None


# ── Presentation file discovery ──────────────────────────────────────────────

def find_presentation_file(method_dir: Path) -> Path | None:
    """
    Find a presentation file (PPTX or PDF) in a method/paper directory.

    Looks for PPTX first (more structured), then PDF.
    Returns the file path or None if not found.
    """
    # Try PPTX first
    pptx_files = list(method_dir.glob("*.pptx"))
    if pptx_files:
        return pptx_files[0]

    # Try PDF
    pdf_files = list(method_dir.glob("*.pdf"))
    if pdf_files:
        return pdf_files[0]

    return None


# ── Paper / method discovery ──────────────────────────────────────────────────

def discover_papers(papers_dir: Path) -> list[str]:
    """
    Return a sorted list of paper names found in papers_dir.
    A paper is valid if papers_dir/{paper}/original.pdf exists.
    """
    if not papers_dir.exists():
        print(f"  [convert] Papers directory not found: {papers_dir}")
        return []
    return sorted(
        p.name for p in papers_dir.iterdir()
        if p.is_dir() and (p / "original.pdf").exists()
    )


def discover_methods(generated_dir: Path) -> list[str]:
    """Return a sorted list of method sub-directories in generated_dir."""
    if not generated_dir.exists():
        print(f"  [convert] Generated samples directory not found: {generated_dir}")
        return []
    return sorted(
        d.name for d in generated_dir.iterdir()
        if d.is_dir() and not d.name.startswith(".")
    )


# ── Per-paper processing ──────────────────────────────────────────────────────

def process_paper(
    paper_name: str,
    papers_dir: Path,
    generated_dir: Path,
    output_dir: Path,
    methods: list[str],
    force: bool = False,
) -> dict:
    """
    Convert the source PDF and all method presentations for one paper.

    Returns a dict with conversion success flags.
    """
    stats: dict = {"paper": paper_name, "orig": False, "methods": {}}
    paper_out = output_dir / paper_name
    paper_out.mkdir(parents=True, exist_ok=True)

    # ── Source paper (orig.md) ─────────────────────────────────────────────
    orig_out = paper_out / "orig.md"
    # Also check for legacy .txt files and migrate them
    orig_out_legacy = paper_out / "orig.txt"
    if orig_out_legacy.exists() and not orig_out.exists():
        orig_out_legacy.rename(orig_out)
        print(f"  [convert] Migrated {paper_name}/orig.txt -> orig.md")

    if orig_out.exists() and not force:
        print(f"  [convert] Skipping {paper_name}/orig.md (already exists)")
        stats["orig"] = True
    else:
        # Try pre-existing source.md first (already converted by other tooling)
        source_md = papers_dir / paper_name / "source.md"
        pdf_path = papers_dir / paper_name / "original.pdf"

        if source_md.exists():
            print(f"  [convert] Using existing source.md for {paper_name}")
            text = source_md.read_text(encoding="utf-8")
            text, was_stripped = strip_supplementary_material(text)
            if was_stripped:
                print(f"    [convert] Stripped supplementary material from {paper_name}")
            orig_out.write_text(text, encoding="utf-8")
            stats["orig"] = True
            print(f"    -> Written {orig_out}")
        elif pdf_path.exists():
            print(f"  [convert] Converting PDF: {paper_name}")
            text = convert_pdf_to_markdown(pdf_path)
            if text:
                text, was_stripped = strip_supplementary_material(text)
                if was_stripped:
                    print(f"    [convert] Stripped supplementary material from {paper_name}")
                orig_out.write_text(text, encoding="utf-8")
                stats["orig"] = True
                print(f"    -> Written {orig_out}")
        else:
            print(f"  [convert] Warning: no source.md or original.pdf found for {paper_name}")

    # ── Method presentations ───────────────────────────────────────────────
    for method in methods:
        method_out = paper_out / f"{method}.md"
        # Also check for legacy .txt files and migrate them
        method_out_legacy = paper_out / f"{method}.txt"
        if method_out_legacy.exists() and not method_out.exists():
            method_out_legacy.rename(method_out)
            print(f"  [convert] Migrated {paper_name}/{method}.txt -> {method}.md")

        if method_out.exists() and not force:
            print(f"  [convert] Skipping {paper_name}/{method}.md (already exists)")
            stats["methods"][method] = True
            continue

        method_dir = generated_dir / method / paper_name
        if not method_dir.exists():
            print(f"  [convert] Warning: directory not found for {method}/{paper_name}")
            stats["methods"][method] = False
            continue

        pres_file = find_presentation_file(method_dir)
        if pres_file is not None:
            print(f"  [convert] Converting {method} / {paper_name} ({pres_file.suffix})")
            if pres_file.suffix.lower() == ".pptx":
                text = convert_pptx_to_text(pres_file)
            else:  # .pdf
                text = convert_pdf_slides_to_text(pres_file)

            # Check if the extraction yielded only empty slides (e.g. image-only PDFs).
            # If so, don't write and fall through to the slide_text.txt fallback.
            all_empty = False
            if text:
                non_empty = [
                    line for line in text.split("\n")
                    if line.strip()
                    and not line.strip().startswith("## Slide")
                    and line.strip() != "(empty slide)"
                ]
                all_empty = len(non_empty) == 0

            if text and not all_empty:
                method_out.write_text(text, encoding="utf-8")
                stats["methods"][method] = True
                print(f"    -> Written {method_out}")
                continue
            elif all_empty:
                print(f"    [convert] PDF/PPTX yielded all empty slides; trying GPT-4o vision extraction")
                # Fallback 1: GPT-4o vision extraction for image-only slides
                gpt4o_text = extract_text_via_gpt4o(pres_file, method, paper_name)
                if gpt4o_text:
                    method_out.write_text(gpt4o_text, encoding="utf-8")
                    stats["methods"][method] = True
                    print(f"    -> Written {method_out} (via GPT-4o)")
                    continue
                print(f"    [convert] GPT-4o extraction failed; trying slide_text.txt fallback")
            else:
                stats["methods"][method] = False
                continue

        # Fallback 2: use existing slide_text.txt if no PPTX/PDF or GPT-4o failed
        slide_text_path = method_dir / "slide_text.txt"
        if slide_text_path.exists():
            print(f"  [convert] Using slide_text.txt for {method} / {paper_name}")
            text = slide_text_path.read_text(encoding="utf-8")
            method_out.write_text(text, encoding="utf-8")
            stats["methods"][method] = True
            print(f"    -> Written {method_out}")
            continue

        # Fallback 3: No PPTX/PDF found at all — try GPT-4o if there are images in the dir
        if C.VISION_EXTRACTION_ENABLED and C.OPENAI_API_KEY:
            from utils.image_utils import find_and_convert_images
            existing_images = find_and_convert_images(method, paper_name)
            if existing_images:
                print(f"  [convert] Found images for {method}/{paper_name}; trying GPT-4o extraction")
                gpt4o_text = _extract_from_existing_images(existing_images)
                if gpt4o_text:
                    method_out.write_text(gpt4o_text, encoding="utf-8")
                    stats["methods"][method] = True
                    print(f"    -> Written {method_out} (via GPT-4o)")
                    continue

        print(f"  [convert] Warning: no PPTX, PDF, or slide_text.txt for {method}/{paper_name}")
        stats["methods"][method] = False

    return stats


# ── Public entry point ────────────────────────────────────────────────────────

def run_conversion(
    papers: list[str] | None = None,
    methods: list[str] | None = None,
    force: bool = False,
) -> list[dict]:
    """
    Convert all papers / methods to markdown files in PROCESSED_DATA_DIR.

    Args:
        papers:  List of paper names to process. None = auto-discover.
        methods: List of method names to process. None = auto-discover.
        force:   Re-convert even if the output file already exists.

    Returns:
        List of per-paper conversion stats dicts.
    """
    papers_dir = Path(C.BENCHMARK_DATA_DIR)
    generated_dir = Path(C.GENERATED_SAMPLES_DIR)
    output_dir = Path(C.PROCESSED_DATA_DIR)
    output_dir.mkdir(parents=True, exist_ok=True)

    if methods is None:
        methods = discover_methods(generated_dir)
    if papers is None:
        papers = discover_papers(papers_dir)

    print(f"[convert] Methods : {methods}")
    print(f"[convert] Papers  : {len(papers)} found")

    all_stats = []
    for i, paper in enumerate(papers, 1):
        print(f"\n[convert] [{i}/{len(papers)}] {paper}")
        stats = process_paper(paper, papers_dir, generated_dir, output_dir, methods, force)
        all_stats.append(stats)

    # Summary
    print("\n" + "=" * 60)
    print("CONVERSION SUMMARY")
    print("=" * 60)
    print(f"Papers with orig.md : {sum(s['orig'] for s in all_stats)}/{len(all_stats)}")
    method_totals: dict[str, int] = {}
    for s in all_stats:
        for m, ok in s["methods"].items():
            method_totals[m] = method_totals.get(m, 0) + ok
    for m in sorted(method_totals):
        print(f"  {m}: {method_totals[m]}/{len(all_stats)}")
    print(f"Output: {output_dir.absolute()}")
    return all_stats
